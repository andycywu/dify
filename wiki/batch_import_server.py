#!/usr/bin/env python3
"""
批量文檔導入到 Wiki.js 的後端處理服務
支援 PDF、PPT、Excel、Word 等格式轉換為 Markdown
"""

import os
import sys
import json
import uuid
import tempfile
import mimetypes
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Optional, Tuple
from dataclasses import dataclass

# 文檔處理相關
import pypandoc
import pypdfium2 as pdfium
from docx import Document as DocxDocument
import openpyxl
from pptx import Presentation
import pandas as pd
import markdown

# Web框架
from flask import Flask, request, jsonify, render_template
from werkzeug.utils import secure_filename

# Wiki.js API 客戶端
import requests
import psycopg2
from psycopg2.extras import RealDictCursor

# 配置
app = Flask(__name__)
app.config['MAX_CONTENT_LENGTH'] = 100 * 1024 * 1024  # 100MB
app.config['UPLOAD_FOLDER'] = '/tmp/wiki_imports'

# 創建臨時目錄
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

@dataclass
class ImportOptions:
    target_folder: str = "/imported"
    page_template: str = "standard"
    naming_rule: str = "original"
    extract_images: bool = True
    preserve_formatting: bool = True
    create_toc: bool = False

@dataclass
class ProcessResult:
    success: bool
    filename: str
    wiki_url: Optional[str] = None
    page_id: Optional[str] = None
    error: Optional[str] = None

class DocumentProcessor:
    """文檔處理器"""
    
    def __init__(self):
        self.supported_formats = {
            'pdf': self._process_pdf,
            'docx': self._process_docx,
            'doc': self._process_doc,
            'xlsx': self._process_xlsx,
            'xls': self._process_excel,
            'pptx': self._process_pptx,
            'ppt': self._process_ppt,
            'txt': self._process_txt,
            'md': self._process_markdown,
            'csv': self._process_csv
        }
    
    def process_file(self, file_path: str, options: ImportOptions) -> Tuple[str, Dict]:
        """處理單個文件並返回 Markdown 內容和元數據"""
        file_ext = Path(file_path).suffix.lower().lstrip('.')
        
        if file_ext not in self.supported_formats:
            raise ValueError(f"不支援的文件格式: {file_ext}")
        
        processor = self.supported_formats[file_ext]
        return processor(file_path, options)
    
    def _process_pdf(self, file_path: str, options: ImportOptions) -> Tuple[str, Dict]:
        """處理 PDF 文件"""
        try:
            pdf = pdfium.PdfDocument(file_path)
            content_parts = []
            
            for page_num in range(len(pdf)):
                page = pdf.get_page(page_num)
                text = page.get_textpage().get_text_range()
                if text.strip():
                    content_parts.append(f"## 第 {page_num + 1} 頁\n\n{text.strip()}")
                page.close()
            
            pdf.close()
            
            content = "\n\n".join(content_parts)
            metadata = {
                'source_type': 'PDF',
                'pages': len(pdf),
                'processed_at': datetime.now().isoformat()
            }
            
            return content, metadata
            
        except Exception as e:
            raise Exception(f"PDF 處理失敗: {str(e)}")
    
    def _process_docx(self, file_path: str, options: ImportOptions) -> Tuple[str, Dict]:
        """處理 DOCX 文件"""
        try:
            doc = DocxDocument(file_path)
            content_parts = []
            
            for paragraph in doc.paragraphs:
                text = paragraph.text.strip()
                if text:
                    # 根據樣式判斷是否為標題
                    if paragraph.style.name.startswith('Heading'):
                        level = int(paragraph.style.name.split()[-1])
                        content_parts.append(f"{'#' * level} {text}")
                    else:
                        content_parts.append(text)
            
            # 處理表格
            for table in doc.tables:
                table_md = self._convert_table_to_markdown(table)
                content_parts.append(table_md)
            
            content = "\n\n".join(content_parts)
            metadata = {
                'source_type': 'DOCX',
                'paragraphs': len(doc.paragraphs),
                'tables': len(doc.tables),
                'processed_at': datetime.now().isoformat()
            }
            
            return content, metadata
            
        except Exception as e:
            raise Exception(f"DOCX 處理失敗: {str(e)}")
    
    def _process_doc(self, file_path: str, options: ImportOptions) -> Tuple[str, Dict]:
        """處理 DOC 文件 (使用 pypandoc)"""
        try:
            # 使用 pypandoc 轉換 DOC 到 Markdown
            content = pypandoc.convert_file(file_path, 'md')
            metadata = {
                'source_type': 'DOC',
                'processed_at': datetime.now().isoformat(),
                'converter': 'pypandoc'
            }
            return content, metadata
        except Exception as e:
            raise Exception(f"DOC 處理失敗: {str(e)}")
    
    def _process_xlsx(self, file_path: str, options: ImportOptions) -> Tuple[str, Dict]:
        """處理 XLSX 文件"""
        try:
            workbook = openpyxl.load_workbook(file_path)
            content_parts = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                content_parts.append(f"## {sheet_name}")
                
                # 轉換為表格格式
                data = []
                for row in sheet.iter_rows(values_only=True):
                    if any(cell is not None for cell in row):
                        data.append([str(cell) if cell is not None else "" for cell in row])
                
                if data:
                    # 創建 Markdown 表格
                    if len(data) > 0:
                        headers = data[0]
                        table_md = "| " + " | ".join(headers) + " |\n"
                        table_md += "| " + " | ".join(["---"] * len(headers)) + " |\n"
                        
                        for row in data[1:]:
                            table_md += "| " + " | ".join(row) + " |\n"
                        
                        content_parts.append(table_md)
            
            content = "\n\n".join(content_parts)
            metadata = {
                'source_type': 'XLSX',
                'sheets': len(workbook.sheetnames),
                'processed_at': datetime.now().isoformat()
            }
            
            return content, metadata
            
        except Exception as e:
            raise Exception(f"XLSX 處理失敗: {str(e)}")
    
    def _process_excel(self, file_path: str, options: ImportOptions) -> Tuple[str, Dict]:
        """處理 XLS 文件"""
        try:
            # 使用 pandas 讀取 XLS
            xls = pd.ExcelFile(file_path)
            content_parts = []
            
            for sheet_name in xls.sheet_names:
                df = pd.read_excel(xls, sheet_name=sheet_name)
                content_parts.append(f"## {sheet_name}")
                
                # 轉換為 Markdown 表格
                table_md = df.to_markdown(index=False)
                content_parts.append(table_md)
            
            content = "\n\n".join(content_parts)
            metadata = {
                'source_type': 'XLS',
                'sheets': len(xls.sheet_names),
                'processed_at': datetime.now().isoformat()
            }
            
            return content, metadata
            
        except Exception as e:
            raise Exception(f"XLS 處理失敗: {str(e)}")
    
    def _process_pptx(self, file_path: str, options: ImportOptions) -> Tuple[str, Dict]:
        """處理 PPTX 文件"""
        try:
            presentation = Presentation(file_path)
            content_parts = []
            
            for i, slide in enumerate(presentation.slides, 1):
                content_parts.append(f"## 投影片 {i}")
                
                # 提取文字內容
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if slide_text:
                    content_parts.append("\n".join(slide_text))
                else:
                    content_parts.append("*(此投影片沒有文字內容)*")
            
            content = "\n\n".join(content_parts)
            metadata = {
                'source_type': 'PPTX',
                'slides': len(presentation.slides),
                'processed_at': datetime.now().isoformat()
            }
            
            return content, metadata
            
        except Exception as e:
            raise Exception(f"PPTX 處理失敗: {str(e)}")
    
    def _process_ppt(self, file_path: str, options: ImportOptions) -> Tuple[str, Dict]:
        """處理 PPT 文件 (使用 pypandoc)"""
        try:
            content = pypandoc.convert_file(file_path, 'md')
            metadata = {
                'source_type': 'PPT',
                'processed_at': datetime.now().isoformat(),
                'converter': 'pypandoc'
            }
            return content, metadata
        except Exception as e:
            raise Exception(f"PPT 處理失敗: {str(e)}")
    
    def _process_txt(self, file_path: str, options: ImportOptions) -> Tuple[str, Dict]:
        """處理 TXT 文件"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            metadata = {
                'source_type': 'TXT',
                'processed_at': datetime.now().isoformat()
            }
            
            return content, metadata
            
        except Exception as e:
            raise Exception(f"TXT 處理失敗: {str(e)}")
    
    def _process_markdown(self, file_path: str, options: ImportOptions) -> Tuple[str, Dict]:
        """處理 Markdown 文件"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            metadata = {
                'source_type': 'Markdown',
                'processed_at': datetime.now().isoformat()
            }
            
            return content, metadata
            
        except Exception as e:
            raise Exception(f"Markdown 處理失敗: {str(e)}")
    
    def _process_csv(self, file_path: str, options: ImportOptions) -> Tuple[str, Dict]:
        """處理 CSV 文件"""
        try:
            df = pd.read_csv(file_path)
            content = df.to_markdown(index=False)
            
            metadata = {
                'source_type': 'CSV',
                'rows': len(df),
                'columns': len(df.columns),
                'processed_at': datetime.now().isoformat()
            }
            
            return content, metadata
            
        except Exception as e:
            raise Exception(f"CSV 處理失敗: {str(e)}")
    
    def _convert_table_to_markdown(self, table) -> str:
        """將 Word 表格轉換為 Markdown"""
        try:
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(cells)
            
            if not rows:
                return ""
            
            # 創建 Markdown 表格
            md_table = "| " + " | ".join(rows[0]) + " |\n"
            md_table += "| " + " | ".join(["---"] * len(rows[0])) + " |\n"
            
            for row in rows[1:]:
                md_table += "| " + " | ".join(row) + " |\n"
            
            return md_table
        except:
            return "*(表格轉換失敗)*"

class WikiJSClient:
    """Wiki.js 客戶端"""
    
    def __init__(self, db_config: Dict):
        self.db_config = db_config
    
    def create_page(self, path: str, title: str, content: str, metadata: Dict) -> str:
        """在 Wiki.js 中創建頁面"""
        try:
            # 連接到 PostgreSQL 數據庫
            with psycopg2.connect(**self.db_config) as conn:
                with conn.cursor(cursor_factory=RealDictCursor) as cursor:
                    # 生成唯一的頁面 ID
                    page_id = str(uuid.uuid4())
                    
                    # 插入頁面記錄
                    cursor.execute("""
                        INSERT INTO pages (id, path, title, description, content, "isPrivate", "isPublished", "contentType", "createdAt", "updatedAt")
                        VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s, %s)
                    """, (
                        page_id,
                        path,
                        title,
                        f"從 {metadata.get('source_type', 'Unknown')} 文件導入",
                        content,
                        False,  # isPrivate
                        True,   # isPublished
                        'markdown',
                        datetime.now(),
                        datetime.now()
                    ))
                    
                    conn.commit()
                    return page_id
                    
        except Exception as e:
            raise Exception(f"創建 Wiki 頁面失敗: {str(e)}")

# 初始化處理器
processor = DocumentProcessor()

# Wiki.js 數據庫配置
WIKI_DB_CONFIG = {
    'host': os.getenv('WIKI_DB_HOST', 'localhost'),
    'port': int(os.getenv('WIKI_DB_PORT', '5432')),
    'database': os.getenv('WIKI_DB_NAME', 'wiki'),
    'user': os.getenv('WIKI_DB_USER', 'postgres'),
    'password': os.getenv('WIKI_DB_PASSWORD', 'difyai123456')
}

wiki_client = WikiJSClient(WIKI_DB_CONFIG)

@app.route('/')
def index():
    """主頁"""
    return render_template_string(open('/Users/andycyw/dify/wiki/batch-document-importer.html').read())

@app.route('/api/wiki/batch-import', methods=['POST'])
def batch_import():
    """批量導入文檔到 Wiki.js"""
    try:
        # 檢查文件
        if 'file' not in request.files:
            return jsonify({'error': '沒有上傳文件'}), 400
        
        file = request.files['file']
        if file.filename == '':
            return jsonify({'error': '未選擇文件'}), 400
        
        # 獲取選項
        options = ImportOptions(
            target_folder=request.form.get('targetFolder', '/imported'),
            page_template=request.form.get('pageTemplate', 'standard'),
            naming_rule=request.form.get('namingRule', 'original'),
            extract_images=request.form.get('extractImages', 'true').lower() == 'true',
            preserve_formatting=request.form.get('preserveFormatting', 'true').lower() == 'true',
            create_toc=request.form.get('createToc', 'false').lower() == 'true'
        )
        
        # 保存臨時文件
        filename = secure_filename(file.filename)
        temp_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{uuid.uuid4()}_{filename}")
        file.save(temp_path)
        
        try:
            # 處理文檔
            content, metadata = processor.process_file(temp_path, options)
            
            # 生成頁面路徑和標題
            base_name = Path(filename).stem
            if options.naming_rule == 'timestamp':
                page_title = f"{base_name}_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
            elif options.naming_rule == 'sequential':
                # 這裡可以實現序號邏輯
                page_title = f"{base_name}_001"
            else:
                page_title = base_name
            
            page_path = f"{options.target_folder.rstrip('/')}/{page_title}".replace('//', '/')
            
            # 添加元數據到內容頂部
            full_content = f"""# {page_title}

> **文檔信息**  
> - 原始文件: {filename}  
> - 文件類型: {metadata.get('source_type', 'Unknown')}  
> - 導入時間: {metadata.get('processed_at', 'Unknown')}  

---

{content}
"""
            
            # 創建 Wiki.js 頁面
            page_id = wiki_client.create_page(page_path, page_title, full_content, metadata)
            
            # 清理臨時文件
            os.unlink(temp_path)
            
            return jsonify({
                'success': True,
                'page_id': page_id,
                'wiki_url': f"/wiki{page_path}",
                'title': page_title,
                'metadata': metadata
            })
            
        except Exception as e:
            # 清理臨時文件
            if os.path.exists(temp_path):
                os.unlink(temp_path)
            raise e
            
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/wiki/supported-formats', methods=['GET'])
def supported_formats():
    """獲取支援的文件格式"""
    return jsonify({
        'formats': list(processor.supported_formats.keys()),
        'descriptions': {
            'pdf': 'PDF 文檔',
            'docx': 'Word 文檔 (新版)',
            'doc': 'Word 文檔 (舊版)',
            'xlsx': 'Excel 電子表格 (新版)',
            'xls': 'Excel 電子表格 (舊版)',
            'pptx': 'PowerPoint 簡報 (新版)',
            'ppt': 'PowerPoint 簡報 (舊版)',
            'txt': '純文字文件',
            'md': 'Markdown 文件',
            'csv': 'CSV 表格文件'
        }
    })

@app.route('/api/wiki/status/<job_id>', methods=['GET'])
def check_status(job_id):
    """檢查處理狀態 (預留用於異步處理)"""
    # 這裡可以實現異步處理狀態檢查
    return jsonify({
        'job_id': job_id,
        'status': 'completed',
        'progress': 100
    })

if __name__ == '__main__':
    print("🚀 Wiki.js 批量文檔導入服務啟動中...")
    print("📚 支援格式:", list(processor.supported_formats.keys()))
    print("🌐 服務地址: http://localhost:5000")
    
    app.run(debug=True, host='0.0.0.0', port=5000)
