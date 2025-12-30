#!/usr/bin/env python3
"""
批量文檔導入到 Wiki.js 的後端處理服務
支援 PDF、PPT、Excel、Word 等格式轉換為 Markdown
"""

import os
import sys
import json
import uuid
import tempfile
import mimetypes
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Optional, Tuple
from dataclasses import dataclass
import logging

# 文檔處理相關
import pypandoc
import pypdfium2 as pdfium
from docx import Document as DocxDocument
import openpyxl
from pptx import Presentation
import pandas as pd
import markdown

# Web框架
from flask import Flask, request, jsonify, render_template, render_template_string, redirect
from werkzeug.utils import secure_filename

# Wiki.js API 客戶端
import requests

# 導入 SMB 同步服務
try:
    from smb_sync_service import SMBSyncService, SyncConfig, FileTracker
    SMB_AVAILABLE = True
except ImportError:
    SMB_AVAILABLE = False
    logging.warning("⚠️ SMB sync service not available")

# 配置
app = Flask(__name__)
app.config['MAX_CONTENT_LENGTH'] = 100 * 1024 * 1024  # 100MB
app.config['UPLOAD_FOLDER'] = '/tmp/wiki_imports'

# 創建臨時目錄
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

@dataclass
class ImportOptions:
    target_folder: str = "/imported"
    page_template: str = "standard"
    naming_rule: str = "original"
    extract_images: bool = True
    preserve_formatting: bool = True
    create_toc: bool = False

@dataclass
class ProcessResult:
    success: bool
    filename: str
    wiki_url: Optional[str] = None
    page_id: Optional[str] = None
    error: Optional[str] = None

class DocumentProcessor:
    """文檔處理器"""

    def __init__(self):
        self.supported_formats = {
            'pdf': self._process_pdf,
            'docx': self._process_docx,
            'doc': self._process_doc,
            'xlsx': self._process_xlsx,
            'xls': self._process_excel,
            'pptx': self._process_pptx,
            'ppt': self._process_ppt,
            'txt': self._process_txt,
            'md': self._process_markdown,
            'csv': self._process_csv
        }

    def process_file(self, file_path: str, options: ImportOptions) -> Tuple[str, Dict]:
        """處理單個文件並返回 Markdown 內容和元數據"""
        file_ext = Path(file_path).suffix.lower().lstrip('.')

        if file_ext not in self.supported_formats:
            raise ValueError(f"不支援的文件格式: {file_ext}")

        processor = self.supported_formats[file_ext]
        return processor(file_path, options)

    def _process_pdf(self, file_path: str, options: ImportOptions) -> Tuple[str, Dict]:
        """處理 PDF 文件"""
        try:
            pdf = pdfium.PdfDocument(file_path)
            content_parts = []

            for page_num in range(len(pdf)):
                page = pdf.get_page(page_num)
                text = page.get_textpage().get_text_range()
                if text.strip():
                    content_parts.append(f"## 第 {page_num + 1} 頁\n\n{text.strip()}")
                page.close()

            pdf.close()

            content = "\n\n".join(content_parts)
            metadata = {
                'source_type': 'PDF',
                'pages': len(pdf),
                'processed_at': datetime.now().isoformat()
            }

            return content, metadata

        except Exception as e:
            raise Exception(f"PDF 處理失敗: {str(e)}")

    def _process_docx(self, file_path: str, options: ImportOptions) -> Tuple[str, Dict]:
        """處理 DOCX 文件"""
        try:
            doc = DocxDocument(file_path)
            content_parts = []

            for paragraph in doc.paragraphs:
                text = paragraph.text.strip()
                if text:
                    # 根據樣式判斷是否為標題
                    if paragraph.style.name.startswith('Heading'):
                        level = int(paragraph.style.name.split()[-1])
                        content_parts.append(f"{'#' * level} {text}")
                    else:
                        content_parts.append(text)

            # 處理表格
            for table in doc.tables:
                table_md = self._convert_table_to_markdown(table)
                content_parts.append(table_md)

            content = "\n\n".join(content_parts)
            metadata = {
                'source_type': 'DOCX',
                'paragraphs': len(doc.paragraphs),
                'tables': len(doc.tables),
                'processed_at': datetime.now().isoformat()
            }

            return content, metadata

        except Exception as e:
            raise Exception(f"DOCX 處理失敗: {str(e)}")

    def _process_doc(self, file_path: str, options: ImportOptions) -> Tuple[str, Dict]:
        """處理 DOC 文件 (使用 pypandoc)"""
        try:
            # 使用 pypandoc 轉換 DOC 到 Markdown
            content = pypandoc.convert_file(file_path, 'md')
            metadata = {
                'source_type': 'DOC',
                'processed_at': datetime.now().isoformat(),
                'converter': 'pypandoc'
            }
            return content, metadata
        except Exception as e:
            raise Exception(f"DOC 處理失敗: {str(e)}")

    def _process_xlsx(self, file_path: str, options: ImportOptions) -> Tuple[str, Dict]:
        """處理 XLSX 文件"""
        try:
            workbook = openpyxl.load_workbook(file_path)
            content_parts = []

            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                content_parts.append(f"## {sheet_name}")

                # 轉換為表格格式
                data = []
                for row in sheet.iter_rows(values_only=True):
                    if any(cell is not None for cell in row):
                        data.append([str(cell) if cell is not None else "" for cell in row])

                if data:
                    # 創建 Markdown 表格
                    if len(data) > 0:
                        headers = data[0]
                        table_md = "| " + " | ".join(headers) + " |\n"
                        table_md += "| " + " | ".join(["---"] * len(headers)) + " |\n"

                        for row in data[1:]:
                            table_md += "| " + " | ".join(row) + " |\n"

                        content_parts.append(table_md)

            content = "\n\n".join(content_parts)
            metadata = {
                'source_type': 'XLSX',
                'sheets': len(workbook.sheetnames),
                'processed_at': datetime.now().isoformat()
            }

            return content, metadata

        except Exception as e:
            raise Exception(f"XLSX 處理失敗: {str(e)}")

    def _process_excel(self, file_path: str, options: ImportOptions) -> Tuple[str, Dict]:
        """處理 XLS 文件"""
        try:
            # 使用 pandas 讀取 XLS
            xls = pd.ExcelFile(file_path)
            content_parts = []

            for sheet_name in xls.sheet_names:
                df = pd.read_excel(xls, sheet_name=sheet_name)
                content_parts.append(f"## {sheet_name}")

                # 轉換為 Markdown 表格
                table_md = df.to_markdown(index=False)
                content_parts.append(table_md)

            content = "\n\n".join(content_parts)
            metadata = {
                'source_type': 'XLS',
                'sheets': len(xls.sheet_names),
                'processed_at': datetime.now().isoformat()
            }

            return content, metadata

        except Exception as e:
            raise Exception(f"XLS 處理失敗: {str(e)}")

    def _process_pptx(self, file_path: str, options: ImportOptions) -> Tuple[str, Dict]:
        """處理 PPTX 文件"""
        try:
            presentation = Presentation(file_path)
            content_parts = []

            for i, slide in enumerate(presentation.slides, 1):
                content_parts.append(f"## 投影片 {i}")

                # 提取文字內容
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())

                if slide_text:
                    content_parts.append("\n".join(slide_text))
                else:
                    content_parts.append("*(此投影片沒有文字內容)*")

            content = "\n\n".join(content_parts)
            metadata = {
                'source_type': 'PPTX',
                'slides': len(presentation.slides),
                'processed_at': datetime.now().isoformat()
            }

            return content, metadata

        except Exception as e:
            raise Exception(f"PPTX 處理失敗: {str(e)}")

    def _process_ppt(self, file_path: str, options: ImportOptions) -> Tuple[str, Dict]:
        """處理 PPT 文件 (使用 pypandoc)"""
        try:
            content = pypandoc.convert_file(file_path, 'md')
            metadata = {
                'source_type': 'PPT',
                'processed_at': datetime.now().isoformat(),
                'converter': 'pypandoc'
            }
            return content, metadata
        except Exception as e:
            raise Exception(f"PPT 處理失敗: {str(e)}")

    def _process_txt(self, file_path: str, options: ImportOptions) -> Tuple[str, Dict]:
        """處理 TXT 文件"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()

            metadata = {
                'source_type': 'TXT',
                'processed_at': datetime.now().isoformat()
            }

            return content, metadata

        except Exception as e:
            raise Exception(f"TXT 處理失敗: {str(e)}")

    def _process_markdown(self, file_path: str, options: ImportOptions) -> Tuple[str, Dict]:
        """處理 Markdown 文件"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()

            metadata = {
                'source_type': 'Markdown',
                'processed_at': datetime.now().isoformat()
            }

            return content, metadata

        except Exception as e:
            raise Exception(f"Markdown 處理失敗: {str(e)}")

    def _process_csv(self, file_path: str, options: ImportOptions) -> Tuple[str, Dict]:
        """處理 CSV 文件"""
        try:
            df = pd.read_csv(file_path)
            content = df.to_markdown(index=False)

            metadata = {
                'source_type': 'CSV',
                'rows': len(df),
                'columns': len(df.columns),
                'processed_at': datetime.now().isoformat()
            }

            return content, metadata

        except Exception as e:
            raise Exception(f"CSV 處理失敗: {str(e)}")

    def _convert_table_to_markdown(self, table) -> str:
        """將 Word 表格轉換為 Markdown"""
        try:
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(cells)

            if not rows:
                return ""

            # 創建 Markdown 表格
            md_table = "| " + " | ".join(rows[0]) + " |\n"
            md_table += "| " + " | ".join(["---"] * len(rows[0])) + " |\n"

            for row in rows[1:]:
                md_table += "| " + " | ".join(row) + " |\n"

            return md_table
        except:
            return "*(表格轉換失敗)*"

class WikiJSClient:
    """Wiki.js 客戶端 - 使用 GraphQL API"""

    def __init__(self, api_url: str = None, api_key: str = None):
        self.api_url = api_url or os.getenv('WIKI_API_URL', 'http://wiki:3000/graphql')
        self.api_key = api_key or os.getenv('WIKI_API_KEY', '')
        self.session = requests.Session()
        if self.api_key:
            self.session.headers.update({'Authorization': f'Bearer {self.api_key}'})

    def create_page(self, path: str, title: str, content: str, metadata: Dict) -> int:
        """使用 GraphQL API 在 Wiki.js 中創建頁面"""
        import sys
        try:
            # 移除開頭的斜線
            clean_path = path.lstrip('/')

            print(f"🔍 Creating page: path={clean_path}, title={title}", file=sys.stderr, flush=True)
            print(f"🔍 API URL: {self.api_url}", file=sys.stderr, flush=True)
            print(f"🔍 Has API Key: {bool(self.api_key)}", file=sys.stderr, flush=True)

            # GraphQL mutation 來創建頁面
            mutation = """
            mutation CreatePage($content: String!, $path: String!, $title: String!) {
              pages {
                create(
                  content: $content
                  description: "從文件導入"
                  editor: "markdown"
                  isPublished: true
                  isPrivate: false
                  locale: "en"
                  path: $path
                  tags: []
                  title: $title
                ) {
                  responseResult {
                    succeeded
                    errorCode
                    slug
                    message
                  }
                  page {
                    id
                    path
                    title
                  }
                }
              }
            }
            """

            variables = {
                "content": content,
                "path": clean_path,
                "title": title
            }

            print(f"🔍 Sending GraphQL request...", file=sys.stderr, flush=True)
            response = self.session.post(
                self.api_url,
                json={"query": mutation, "variables": variables},
                timeout=30
            )

            print(f"🔍 Response status: {response.status_code}", file=sys.stderr, flush=True)
            print(f"🔍 Response body: {response.text[:500]}", file=sys.stderr, flush=True)

            if response.status_code == 200:
                result = response.json()

                if 'errors' in result:
                    # 如果是因為頁面已存在，嘗試更新
                    error_msg = str(result['errors'])
                    print(f"❌ GraphQL errors: {error_msg}", file=sys.stderr, flush=True)
                    if 'already exists' in error_msg.lower() or 'duplicate' in error_msg.lower():
                        return self.update_page(clean_path, title, content, metadata)
                    raise Exception(f"GraphQL 錯誤: {result['errors']}")

                page_result = result.get('data', {}).get('pages', {}).get('create', {})
                response_result = page_result.get('responseResult', {})

                print(f"🔍 Response result: {response_result}", file=sys.stderr, flush=True)

                if response_result.get('succeeded'):
                    page_info = page_result.get('page', {})
                    page_id = page_info.get('id', 0)
                    print(f"✅ Page created successfully with ID: {page_id}", file=sys.stderr, flush=True)
                    return page_id
                else:
                    error_msg = response_result.get('message', 'Unknown error')
                    print(f"❌ Page creation failed: {error_msg}", file=sys.stderr, flush=True)
                    # 如果是重複頁面錯誤，嘗試更新
                    if 'already exists' in error_msg.lower():
                        return self.update_page(clean_path, title, content, metadata)
                    raise Exception(f"創建頁面失敗: {error_msg}")
            else:
                print(f"❌ HTTP error: {response.status_code}", file=sys.stderr, flush=True)
                raise Exception(f"HTTP 錯誤: {response.status_code}")

        except Exception as e:
            print(f"❌ Exception in create_page: {str(e)}", file=sys.stderr, flush=True)
            raise Exception(f"創建 Wiki 頁面失敗: {str(e)}")

    def update_page(self, path: str, title: str, content: str, metadata: Dict) -> int:
        """更新現有頁面"""
        try:
            clean_path = path.lstrip('/')

            # 首先獲取頁面 ID
            query = """
            query GetPage($path: String!) {
              pages {
                single(path: $path) {
                  id
                }
              }
            }
            """

            response = self.session.post(
                self.api_url,
                json={"query": query, "variables": {"path": clean_path}},
                timeout=10
            )

            if response.status_code == 200:
                result = response.json()
                page_id = result.get('data', {}).get('pages', {}).get('single', {}).get('id')

                if not page_id:
                    raise Exception("找不到頁面進行更新")

                # 更新頁面
                mutation = """
                mutation UpdatePage($id: Int!, $content: String!, $title: String!) {
                  pages {
                    update(
                      id: $id
                      content: $content
                      title: $title
                      description: "從文件導入（已更新）"
                      editor: "markdown"
                      isPublished: true
                      isPrivate: false
                      locale: "en"
                    ) {
                      responseResult {
                        succeeded
                        errorCode
                        message
                      }
                    }
                  }
                }
                """

                variables = {
                    "id": page_id,
                    "content": content,
                    "title": title
                }

                response = self.session.post(
                    self.api_url,
                    json={"query": mutation, "variables": variables},
                    timeout=30
                )

                if response.status_code == 200:
                    result = response.json()
                    update_result = result.get('data', {}).get('pages', {}).get('update', {}).get('responseResult', {})

                    if update_result.get('succeeded'):
                        return page_id
                    else:
                        raise Exception(f"更新失敗: {update_result.get('message')}")
                else:
                    raise Exception(f"HTTP 錯誤: {response.status_code}")
            else:
                raise Exception(f"HTTP 錯誤: {response.status_code}")

        except Exception as e:
            raise Exception(f"更新 Wiki 頁面失敗: {str(e)}")

# 初始化處理器
processor = DocumentProcessor()

# Wiki.js API 配置
WIKI_API_URL = os.getenv('WIKI_API_URL', 'http://wiki:3000/graphql')
WIKI_API_KEY = os.getenv('WIKI_API_KEY', '')

wiki_client = WikiJSClient(api_url=WIKI_API_URL, api_key=WIKI_API_KEY)

# ========== API 路由 ==========

@app.route('/api/wiki/batch-import', methods=['POST'])
def batch_import():
    """批量導入文檔到 Wiki.js"""
    try:
        # 檢查文件
        if 'file' not in request.files:
            return jsonify({'error': '沒有上傳文件'}), 400

        file = request.files['file']
        if file.filename == '':
            return jsonify({'error': '未選擇文件'}), 400

        # 獲取選項
        options = ImportOptions(
            target_folder=request.form.get('targetFolder', '/imported'),
            page_template=request.form.get('pageTemplate', 'standard'),
            naming_rule=request.form.get('namingRule', 'original'),
            extract_images=request.form.get('extractImages', 'true').lower() == 'true',
            preserve_formatting=request.form.get('preserveFormatting', 'true').lower() == 'true',
            create_toc=request.form.get('createToc', 'false').lower() == 'true'
        )

        # 保存臨時文件
        filename = secure_filename(file.filename)
        temp_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{uuid.uuid4()}_{filename}")
        file.save(temp_path)

        try:
            # 處理文檔
            content, metadata = processor.process_file(temp_path, options)

            # 生成頁面路徑和標題
            base_name = Path(filename).stem
            if options.naming_rule == 'timestamp':
                page_title = f"{base_name}_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
            elif options.naming_rule == 'sequential':
                # 這裡可以實現序號邏輯
                page_title = f"{base_name}_001"
            else:
                page_title = base_name

            page_path = f"{options.target_folder.rstrip('/')}/{page_title}".replace('//', '/')

            # 添加元數據到內容頂部
            full_content = f"""# {page_title}

> **文檔信息**
> - 原始文件: {filename}
> - 文件類型: {metadata.get('source_type', 'Unknown')}
> - 導入時間: {metadata.get('processed_at', 'Unknown')}

---

{content}
"""

            # 創建 Wiki.js 頁面
            page_id = wiki_client.create_page(page_path, page_title, full_content, metadata)

            # 清理臨時文件
            os.unlink(temp_path)

            return jsonify({
                'success': True,
                'page_id': page_id,
                'wiki_url': f"/wiki{page_path}",
                'title': page_title,
                'metadata': metadata
            })

        except Exception as e:
            # 清理臨時文件
            if os.path.exists(temp_path):
                os.unlink(temp_path)
            raise e

    except Exception as e:
        return jsonify({'error': str(e)}), 500

# ========== Web 管理界面路由 ==========

@app.route('/')
def index():
    """主頁面 - Web 管理界面"""
    # 檢查 URL 參數中的認證 token
    auth_token = request.args.get('auth')
    print(f"[DEBUG] Received auth_token: {auth_token is not None}")
    print(f"[DEBUG] Full URL: {request.url}")
    print(f"[DEBUG] Request args: {dict(request.args)}")

    if auth_token:
        try:
            # 驗證 token
            import base64
            import json
            import time

            print(f"[DEBUG] Decoding token: {auth_token[:50]}...")
            token_data = json.loads(base64.b64decode(auth_token).decode('utf-8'))
            print(f"[DEBUG] Decoded token data: {token_data}")

            # 檢查過期時間
            if token_data.get('exp', 0) < time.time():
                print(f"[DEBUG] Token expired: {token_data.get('exp')} < {time.time()}")
                return '''
                <html>
                    <body>
                        <h1>認證過期</h1>
                        <p>請重新從主系統登入。</p>
                        <a href="''' + ADMIN_FRONTEND_URL + '''">返回主系統</a>
                    </body>
                </html>
                '''

            # 檢查角色
            user_role = token_data.get('role')
            print(f"[DEBUG] Checking role: '{user_role}' (type: {type(user_role)})")
            if user_role != 'Administrator':
                print(f"[DEBUG] Invalid role: '{user_role}' != 'Administrator'")
                return '''
                <html>
                    <body>
                        <h1>權限不足</h1>
                        <p>需要管理員權限。當前角色: {user_role}</p>
                        <a href="''' + ADMIN_FRONTEND_URL + '''">返回主系統</a>
                    </body>
                </html>
                '''.format(user_role=user_role)

            # Token 有效，重定向到管理界面
            print(f"[DEBUG] Token valid, redirecting to /admin")
            return redirect('/admin')

        except Exception as e:
            print(f"[DEBUG] Token validation error: {str(e)}")
            return f'''
            <html>
                <body>
                    <h1>認證錯誤</h1>
                    <p>無效的認證 token: {str(e)}</p>
                    <a href="''' + ADMIN_FRONTEND_URL + '''">返回主系統</a>
                </body>
            </html>
            '''

    # 沒有認證 token，顯示登入提示
    return '''
    <html>
        <head>
            <title>Wiki Batch Importer</title>
            <style>
                body { font-family: Arial, sans-serif; margin: 40px; background-color: #f5f5f5; }
                .container { max-width: 800px; margin: 0 auto; background: white; padding: 30px; border-radius: 8px; box-shadow: 0 2px 10px rgba(0,0,0,0.1); }
                .warning { background: #fff3cd; border: 1px solid #ffeaa7; color: #856404; padding: 15px; border-radius: 4px; margin-bottom: 20px; }
                .button { display: inline-block; padding: 10px 20px; background: #007bff; color: white; text-decoration: none; border-radius: 4px; margin: 5px; }
                .button:hover { background: #0056b3; }
            </style>
        </head>
        <body>
            <div class="container">
                <h1>📂 Wiki Batch Importer</h1>
                <div class="warning">
                    <strong>🔐 認證要求</strong><br>
                    此工具僅供管理員使用。您需要先從主系統登入才能訪問完整功能。
                </div>
                <p>此工具用於批量導入文檔到 Wiki.js 知識庫系統。</p>

                <h3>📋 支援的功能</h3>
                <ul>
                    <li>批量文檔上傳和導入</li>
                    <li>目錄掃描和自動導入</li>
                    <li>SMB 網路共享同步</li>
                    <li>導入進度追蹤</li>
                </ul>

                <h3>🚀 開始使用</h3>
                <ol>
                    <li>確保您已登入主系統並具有管理員權限</li>
                    <li>使用上方導航返回主系統</li>
                    <li>從管理面板訪問此工具</li>
                </ol>

                <div style="margin-top: 30px;">
                    <a href="''' + ADMIN_FRONTEND_URL + '''" class="button">← 返回主系統</a>
                </div>
            </div>
        </body>
    </html>
    '''

@app.route('/admin')
def admin_panel():
    """管理員界面 - 批量導入工具"""
    return '''
    <html>
        <head>
            <title>Wiki Batch Importer - 管理面板</title>
            <style>
                body { font-family: Arial, sans-serif; margin: 0; padding: 0; background-color: #f8f9fa; }
                .header { background: #343a40; color: white; padding: 15px 20px; display: flex; justify-content: space-between; align-items: center; }
                .header h1 { margin: 0; font-size: 24px; }
                .logout-btn { background: #dc3545; color: white; padding: 8px 16px; text-decoration: none; border-radius: 4px; }
                .logout-btn:hover { background: #c82333; }
                .container { max-width: 1200px; margin: 20px auto; padding: 0 20px; }
                .card { background: white; border-radius: 8px; box-shadow: 0 2px 10px rgba(0,0,0,0.1); margin-bottom: 20px; }
                .card-header { background: #007bff; color: white; padding: 15px 20px; border-radius: 8px 8px 0 0; }
                .card-body { padding: 20px; }
                .form-group { margin-bottom: 15px; }
                .form-group label { display: block; margin-bottom: 5px; font-weight: bold; }
                .form-control { width: 100%; padding: 8px; border: 1px solid #ddd; border-radius: 4px; }
                .btn { padding: 10px 20px; border: none; border-radius: 4px; cursor: pointer; text-decoration: none; display: inline-block; margin: 5px; }
                .btn-primary { background: #007bff; color: white; }
                .btn-primary:hover { background: #0056b3; }
                .btn-success { background: #28a745; color: white; }
                .btn-success:hover { background: #218838; }
                .btn-info { background: #17a2b8; color: white; }
                .btn-info:hover { background: #138496; }
                .status { padding: 10px; border-radius: 4px; margin: 10px 0; }
                .status.success { background: #d4edda; color: #155724; border: 1px solid #c3e6cb; }
                .status.error { background: #f8d7da; color: #721c24; border: 1px solid #f5c6cb; }
                .status.info { background: #d1ecf1; color: #0c5460; border: 1px solid #bee5eb; }
                .progress { width: 100%; height: 20px; background: #f0f0f0; border-radius: 10px; overflow: hidden; margin: 10px 0; }
                .progress-bar { height: 100%; background: #007bff; width: 0%; transition: width 0.3s; }
                .table { width: 100%; border-collapse: collapse; margin-top: 20px; }
                .table th, .table td { padding: 12px; text-align: left; border-bottom: 1px solid #ddd; }
                .table th { background: #f8f9fa; font-weight: bold; }
                .tabs { display: flex; margin-bottom: 20px; }
                .tab { padding: 10px 20px; background: #e9ecef; border: none; cursor: pointer; margin-right: 5px; border-radius: 4px 4px 0 0; }
                .tab.active { background: white; border-bottom: 2px solid #007bff; }
                .tab-content { display: none; }
                .tab-content.active { display: block; }
            </style>
            <script>
                let currentTab = 'upload';
                let currentJobId = null;

                function switchTab(tabName) {
                    document.querySelectorAll('.tab').forEach(t => t.classList.remove('active'));
                    document.querySelectorAll('.tab-content').forEach(c => c.classList.remove('active'));
                    document.getElementById('tab-' + tabName).classList.add('active');
                    document.getElementById(tabName).classList.add('active');
                    currentTab = tabName;
                }

                function uploadFiles() {
                    const files = document.getElementById('files').files;
                    if (files.length === 0) {
                        alert('請選擇要上傳的文件');
                        return;
                    }

                    const formData = new FormData();
                    for (let file of files) {
                        formData.append('files', file);
                    }

                    const statusDiv = document.getElementById('upload-status');
                    statusDiv.innerHTML = '<div class="status info">正在上傳文件...</div>';

                    fetch('/api/wiki/batch-import', {
                        method: 'POST',
                        body: formData
                    })
                    .then(response => response.json())
                    .then(data => {
                        if (data.success) {
                            statusDiv.innerHTML = '<div class="status success">文件上傳成功！任務 ID: ' + data.job_id + '</div>';
                            currentJobId = data.job_id;
                            checkStatus();
                        } else {
                            statusDiv.innerHTML = '<div class="status error">上傳失敗: ' + data.error + '</div>';
                        }
                    })
                    .catch(error => {
                        statusDiv.innerHTML = '<div class="status error">上傳錯誤: ' + error.message + '</div>';
                    });
                }

                function scanDirectory() {
                    const path = document.getElementById('scan-path').value;
                    if (!path) {
                        alert('請輸入目錄路徑');
                        return;
                    }

                    const statusDiv = document.getElementById('scan-status');
                    statusDiv.innerHTML = '<div class="status info">正在掃描目錄...</div>';

                    fetch('/api/wiki/scan-directory', {
                        method: 'POST',
                        headers: { 'Content-Type': 'application/json' },
                        body: JSON.stringify({ path: path })
                    })
                    .then(response => response.json())
                    .then(data => {
                        if (data.success) {
                            statusDiv.innerHTML = '<div class="status success">掃描完成！找到 ' + data.files.length + ' 個文件</div>';
                            displayScanResults(data.files);
                        } else {
                            statusDiv.innerHTML = '<div class="status error">掃描失敗: ' + data.error + '</div>';
                        }
                    })
                    .catch(error => {
                        statusDiv.innerHTML = '<div class="status error">掃描錯誤: ' + error.message + '</div>';
                    });
                }

                function displayScanResults(files) {
                    const tbody = document.getElementById('scan-results');
                    tbody.innerHTML = '';
                    files.forEach(file => {
                        const row = tbody.insertRow();
                        row.insertCell(0).textContent = file.name;
                        row.insertCell(1).textContent = file.path;
                        row.insertCell(2).textContent = (file.size / 1024 / 1024).toFixed(2) + ' MB';
                        row.insertCell(3).textContent = file.type;
                    });
                }

                function importScannedFiles() {
                    const selectedFiles = [];
                    document.querySelectorAll('#scan-results input[type="checkbox"]:checked').forEach(cb => {
                        selectedFiles.push(cb.value);
                    });

                    if (selectedFiles.length === 0) {
                        alert('請選擇要導入的文件');
                        return;
                    }

                    const statusDiv = document.getElementById('scan-status');
                    statusDiv.innerHTML = '<div class="status info">正在導入文件...</div>';

                    fetch('/api/wiki/batch-directory-import', {
                        method: 'POST',
                        headers: { 'Content-Type': 'application/json' },
                        body: JSON.stringify({ files: selectedFiles })
                    })
                    .then(response => response.json())
                    .then(data => {
                        if (data.success) {
                            statusDiv.innerHTML = '<div class="status success">導入任務已啟動！任務 ID: ' + data.job_id + '</div>';
                            currentJobId = data.job_id;
                            checkStatus();
                        } else {
                            statusDiv.innerHTML = '<div class="status error">導入失敗: ' + data.error + '</div>';
                        }
                    })
                    .catch(error => {
                        statusDiv.innerHTML = '<div class="status error">導入錯誤: ' + error.message + '</div>';
                    });
                }

                function checkStatus() {
                    if (!currentJobId) return;

                    fetch('/api/wiki/status/' + currentJobId)
                    .then(response => response.json())
                    .then(data => {
                        const progressBar = document.getElementById('progress-bar');
                        const statusDiv = document.getElementById('job-status');

                        if (progressBar) {
                            progressBar.style.width = data.progress + '%';
                        }

                        if (data.status === 'completed') {
                            statusDiv.innerHTML = '<div class="status success">任務完成！處理了 ' + data.processed + ' 個文件</div>';
                        } else if (data.status === 'failed') {
                            statusDiv.innerHTML = '<div class="status error">任務失敗: ' + data.error + '</div>';
                        } else {
                            statusDiv.innerHTML = '<div class="status info">任務進行中... ' + data.progress + '% (' + data.processed + '/' + data.total + ')</div>';
                            setTimeout(checkStatus, 2000);
                        }
                    })
                    .catch(error => {
                        console.error('檢查狀態失敗:', error);
                    });
                }

                function syncSMB() {
                    const configId = document.getElementById('smb-config').value;
                    if (!configId) {
                        alert('請選擇 SMB 配置');
                        return;
                    }

                    const statusDiv = document.getElementById('smb-status');
                    statusDiv.innerHTML = '<div class="status info">正在同步 SMB 共享...</div>';

                    fetch('/api/wiki/smb-sync', {
                        method: 'POST',
                        headers: { 'Content-Type': 'application/json' },
                        body: JSON.stringify({ config_id: configId })
                    })
                    .then(response => response.json())
                    .then(data => {
                        if (data.success) {
                            statusDiv.innerHTML = '<div class="status success">SMB 同步任務已啟動！任務 ID: ' + data.job_id + '</div>';
                            currentJobId = data.job_id;
                            checkSMBStatus();
                        } else {
                            statusDiv.innerHTML = '<div class="status error">同步失敗: ' + data.error + '</div>';
                        }
                    })
                    .catch(error => {
                        statusDiv.innerHTML = '<div class="status error">同步錯誤: ' + error.message + '</div>';
                    });
                }

                function checkSMBStatus() {
                    if (!currentJobId) return;

                    fetch('/api/wiki/smb-status')
                    .then(response => response.json())
                    .then(data => {
                        const statusDiv = document.getElementById('smb-status');

                        if (data.status === 'running') {
                            statusDiv.innerHTML = '<div class="status info">SMB 同步進行中... ' + data.message + '</div>';
                            setTimeout(checkSMBStatus, 3000);
                        } else if (data.status === 'completed') {
                            statusDiv.innerHTML = '<div class="status success">SMB 同步完成！</div>';
                        } else {
                            statusDiv.innerHTML = '<div class="status error">SMB 同步失敗: ' + data.error + '</div>';
                        }
                    })
                    .catch(error => {
                        console.error('檢查 SMB 狀態失敗:', error);
                    });
                }

                // 初始化
                document.addEventListener('DOMContentLoaded', function() {
                    switchTab('upload');
                });
            </script>
        </head>
        <body>
            <div class="header">
                <h1>📂 Wiki Batch Importer - 管理面板</h1>
                <a href="''' + ADMIN_FRONTEND_URL + '''" class="logout-btn">← 返回主系統</a>
            </div>

            <div class="container">
                <div class="tabs">
                    <button id="tab-upload" class="tab active" onclick="switchTab('upload')">文件上傳</button>
                    <button id="tab-directory" class="tab" onclick="switchTab('directory')">目錄掃描</button>
                    <button id="tab-smb" class="tab" onclick="switchTab('smb')">SMB 同步</button>
                </div>

                <div id="upload" class="tab-content active">
                    <div class="card">
                        <div class="card-header">
                            <h3>📤 批量文件上傳</h3>
                        </div>
                        <div class="card-body">
                            <div class="form-group">
                                <label for="files">選擇文件 (支援 PDF、Word、Excel、PowerPoint 等格式)</label>
                                <input type="file" id="files" multiple class="form-control" accept=".pdf,.doc,.docx,.xls,.xlsx,.ppt,.pptx,.txt,.md">
                            </div>
                            <button onclick="uploadFiles()" class="btn btn-primary">開始上傳</button>
                            <div id="upload-status"></div>
                        </div>
                    </div>
                </div>

                <div id="directory" class="tab-content">
                    <div class="card">
                        <div class="card-header">
                            <h3>📁 目錄掃描與導入</h3>
                        </div>
                        <div class="card-body">
                            <div class="form-group">
                                <label for="scan-path">目錄路徑</label>
                                <input type="text" id="scan-path" class="form-control" placeholder="/path/to/directory">
                            </div>
                            <button onclick="scanDirectory()" class="btn btn-info">掃描目錄</button>
                            <div id="scan-status"></div>

                            <div id="scan-results-container" style="display: none;">
                                <h4>掃描結果</h4>
                                <table class="table">
                                    <thead>
                                        <tr>
                                            <th>文件名</th>
                                            <th>路徑</th>
                                            <th>大小</th>
                                            <th>類型</th>
                                            <th>選擇</th>
                                        </tr>
                                    </thead>
                                    <tbody id="scan-results">
                                    </tbody>
                                </table>
                                <button onclick="importScannedFiles()" class="btn btn-success">導入選中文件</button>
                            </div>
                        </div>
                    </div>
                </div>

                <div id="smb" class="tab-content">
                    <div class="card">
                        <div class="card-header">
                            <h3>🌐 SMB 網路共享同步</h3>
                        </div>
                        <div class="card-body">
                            <div class="form-group">
                                <label for="smb-config">SMB 配置</label>
                                <select id="smb-config" class="form-control">
                                    <option value="">選擇配置...</option>
                                </select>
                            </div>
                            <button onclick="syncSMB()" class="btn btn-success">開始同步</button>
                            <div id="smb-status"></div>
                        </div>
                    </div>
                </div>

                <div id="job-status" class="card">
                    <div class="card-header">
                        <h3>📊 任務狀態</h3>
                    </div>
                    <div class="card-body">
                        <div class="progress">
                            <div id="progress-bar" class="progress-bar"></div>
                        </div>
                        <div id="status-text">等待任務開始...</div>
                    </div>
                </div>
            </div>
        </body>
    </html>
    '''

@app.route('/api/smb-configs', methods=['GET'])
def get_smb_configs():
    """獲取 SMB 配置列表"""
    # TODO: 從數據庫或配置文件讀取
    configs = []
    config_file = '/app/data/smb_configs.json'

    if os.path.exists(config_file):
        try:
            with open(config_file, 'r', encoding='utf-8') as f:
                configs = json.load(f)
        except Exception as e:
            print(f"讀取配置失敗: {e}")

    return jsonify({'success': True, 'configs': configs})

@app.route('/api/smb-configs', methods=['POST'])
def save_smb_config():
    """保存 SMB 配置"""
    try:
        config = request.get_json()
        config_file = '/app/data/smb_configs.json'

        # 創建目錄
        os.makedirs('/app/data', exist_ok=True)

        # 讀取現有配置
        configs = []
        if os.path.exists(config_file):
            with open(config_file, 'r', encoding='utf-8') as f:
                configs = json.load(f)

        # 添加或更新配置
        existing = next((c for c in configs if c['group'] == config['group']), None)
        if existing:
            configs[configs.index(existing)] = config
        else:
            configs.append(config)

        # 保存配置
        with open(config_file, 'w', encoding='utf-8') as f:
            json.dump(configs, f, ensure_ascii=False, indent=2)

        return jsonify({'success': True, 'message': '配置已保存'})
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/api/wiki/supported-formats', methods=['GET'])
def supported_formats():
    """獲取支援的文件格式"""
    return jsonify({
        'formats': list(processor.supported_formats.keys()),
        'descriptions': {
            'pdf': 'PDF 文檔',
            'docx': 'Word 文檔 (新版)',
            'doc': 'Word 文檔 (舊版)',
            'xlsx': 'Excel 電子表格 (新版)',
            'xls': 'Excel 電子表格 (舊版)',
            'pptx': 'PowerPoint 簡報 (新版)',
            'ppt': 'PowerPoint 簡報 (舊版)',
            'txt': '純文字文件',
            'md': 'Markdown 文件',
            'csv': 'CSV 表格文件'
        }
    })

@app.route('/api/wiki/status/<job_id>', methods=['GET'])
def check_status(job_id):
    """檢查處理狀態 (預留用於異步處理)"""
    # 這裡可以實現異步處理狀態檢查
    return jsonify({
        'job_id': job_id,
        'status': 'completed',
        'progress': 100
    })

@app.route('/api/wiki/batch-directory-import', methods=['POST'])
def batch_directory_import():
    """批量目錄導入到 Wiki.js"""
    try:
        # 獲取參數
        data = request.get_json() if request.is_json else request.form
        source_path = data.get('sourcePath', '')
        target_folder = data.get('targetFolder', '/imported')
        preserve_structure = data.get('preserveStructure', 'true').lower() == 'true'

        if not source_path or not os.path.exists(source_path):
            return jsonify({'error': '源路徑不存在'}), 400

        # 掃描目錄
        files_to_import = []
        for root, dirs, files in os.walk(source_path):
            # 排除隱藏目錄
            dirs[:] = [d for d in dirs if not d.startswith('.')]

            for filename in files:
                # 跳過隱藏文件
                if filename.startswith('.'):
                    continue

                # 檢查文件擴展名
                ext = Path(filename).suffix.lower().lstrip('.')
                if ext not in processor.supported_formats:
                    continue

                file_path = os.path.join(root, filename)
                rel_path = os.path.relpath(file_path, source_path)
                files_to_import.append((file_path, rel_path))

        if not files_to_import:
            return jsonify({'error': '目錄中沒有找到支持的文件'}), 400

        # 處理選項
        options = ImportOptions(
            target_folder=target_folder,
            page_template='standard',
            naming_rule='original',
            extract_images=True,
            preserve_formatting=True,
            create_toc=True
        )

        # 批量導入
        results = []
        success_count = 0
        failed_count = 0

        for file_path, rel_path in files_to_import:
            try:
                # 生成 Wiki 路徑
                if preserve_structure:
                    # 保留目錄結構
                    rel_dir = os.path.dirname(rel_path)
                    file_name = Path(file_path).stem
                    if rel_dir:
                        page_path = f"{target_folder}/{rel_dir}/{file_name}".replace('\\', '/').replace('//', '/')
                    else:
                        page_path = f"{target_folder}/{file_name}"
                else:
                    # 扁平化結構
                    file_name = Path(file_path).stem
                    page_path = f"{target_folder}/{file_name}"

                # 處理文檔
                content, metadata = processor.process_file(file_path, options)

                # 添加文件信息
                full_content = f"""# {file_name}

> **文件信息**
> - 原始路徑: `{rel_path}`
> - 文件類型: {metadata.get('source_type', 'Unknown')}
> - 導入時間: {metadata.get('processed_at', 'Unknown')}

---

{content}
"""

                # 創建頁面
                page_id = wiki_client.create_page(page_path, file_name, full_content, metadata)

                results.append({
                    'file': rel_path,
                    'success': True,
                    'page_id': page_id,
                    'wiki_url': f"/wiki{page_path}"
                })
                success_count += 1

            except Exception as e:
                results.append({
                    'file': rel_path,
                    'success': False,
                    'error': str(e)
                })
                failed_count += 1

        return jsonify({
            'success': True,
            'total': len(files_to_import),
            'success_count': success_count,
            'failed_count': failed_count,
            'results': results
        })

    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/wiki/scan-directory', methods=['POST'])
def scan_directory():
    """掃描目錄並返回文件列表"""
    try:
        data = request.get_json() if request.is_json else request.form
        source_path = data.get('sourcePath', '')

        if not source_path or not os.path.exists(source_path):
            return jsonify({'error': '路徑不存在'}), 400

        files = []
        for root, dirs, filenames in os.walk(source_path):
            dirs[:] = [d for d in dirs if not d.startswith('.')]

            for filename in filenames:
                if filename.startswith('.'):
                    continue

                ext = Path(filename).suffix.lower().lstrip('.')
                if ext not in processor.supported_formats:
                    continue

                file_path = os.path.join(root, filename)
                rel_path = os.path.relpath(file_path, source_path)
                file_size = os.path.getsize(file_path)

                files.append({
                    'name': filename,
                    'path': rel_path,
                    'size': file_size,
                    'type': ext
                })

        return jsonify({
            'success': True,
            'path': source_path,
            'total': len(files),
            'files': files
        })

    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/wiki/smb-sync', methods=['POST'])
def smb_sync():
    """SMB 同步服務"""
    if not SMB_AVAILABLE:
        return jsonify({'error': 'SMB sync service not available'}), 503

    try:
        data = request.get_json() if request.is_json else request.form

        # 獲取參數
        group = data.get('group', '')
        smb_path = data.get('smbPath', '')
        target_folder = data.get('targetFolder', f'/smb/{group}')
        mode = data.get('mode', 'once')  # 'once' or 'continuous'
        scan_interval = int(data.get('scanInterval', 300))

        if not group or not smb_path:
            return jsonify({'error': '缺少必要參數: group, smbPath'}), 400

        if not os.path.exists(smb_path):
            return jsonify({'error': f'SMB 路徑不存在: {smb_path}'}), 400

        # 創建同步配置
        config = SyncConfig(
            smb_mount_point=smb_path,
            wiki_base_path=target_folder,
            group_name=group,
            scan_interval=scan_interval,
            state_file=f"/app/data/sync_state_{group}.json"
        )

        # 初始化同步服務
        sync_service = SMBSyncService(
            config=config,
            wiki_client=wiki_client,
            processor=processor
        )

        if mode == 'once':
            # 單次同步
            stats = sync_service.sync_all()

            return jsonify({
                'success': True,
                'mode': 'once',
                'group': group,
                'stats': stats
            })

        elif mode == 'continuous':
            # 連續監控 (在後台線程中運行)
            import threading

            def run_continuous_sync():
                sync_service.start_monitoring()

            sync_thread = threading.Thread(target=run_continuous_sync, daemon=True)
            sync_thread.start()

            return jsonify({
                'success': True,
                'mode': 'continuous',
                'group': group,
                'message': f'已啟動連續監控，掃描間隔 {scan_interval} 秒'
            })

        else:
            return jsonify({'error': f'不支持的模式: {mode}'}), 400

    except Exception as e:
        import traceback
        return jsonify({
            'error': str(e),
            'traceback': traceback.format_exc()
        }), 500

@app.route('/api/wiki/smb-status', methods=['GET'])
def smb_status():
    """檢查 SMB 同步狀態"""
    if not SMB_AVAILABLE:
        return jsonify({'error': 'SMB sync service not available'}), 503

    try:
        group = request.args.get('group', '')

        if not group:
            # 返回所有可用組的狀態
            state_files = Path('/app/data').glob('sync_state_*.json')
            groups = []

            for state_file in state_files:
                group_name = state_file.stem.replace('sync_state_', '')
                tracker = FileTracker(str(state_file))

                groups.append({
                    'group': group_name,
                    'file_count': len(tracker.state['files']),
                    'last_scan': tracker.state.get('last_scan', 'Never')
                })

            return jsonify({
                'success': True,
                'groups': groups
            })

        else:
            # 返回特定組的狀態
            state_file = f"/app/data/sync_state_{group}.json"

            if not os.path.exists(state_file):
                return jsonify({'error': f'找不到組 {group} 的狀態文件'}), 404

            tracker = FileTracker(state_file)

            return jsonify({
                'success': True,
                'group': group,
                'file_count': len(tracker.state['files']),
                'last_scan': tracker.state.get('last_scan', 'Never'),
                'files': tracker.state['files']
            })

    except Exception as e:
        return jsonify({'error': str(e)}), 500

# 使用環境變數配置 API 端點，預設使用 Docker 內部網路
ADMIN_API_VALIDATE = os.getenv("ADMIN_API_VALIDATE", "http://dify-next-frontend:3000/api/auth/validate")
ADMIN_FRONTEND_URL = os.getenv("ADMIN_FRONTEND_URL", "http://localhost:3001")

@app.before_request
def authenticate():
    print(f"[DEBUG] Before request: {request.method} {request.path}")
    print(f"[DEBUG] Request args: {dict(request.args)}")

    # 允許訪問主頁面，用於顯示登入提示
    if request.path == '/':
        print("[DEBUG] Allowing access to root path")
        return

    # 允許訪問管理界面 (認證已在主頁面處理)
    if request.path == '/admin':
        print("[DEBUG] Allowing access to admin path")
        return

    # 檢查 URL 參數中的認證 token
    auth_token = request.args.get('auth')
    if auth_token:
        try:
            # 解碼 token
            import base64
            import json
            import time

            print(f"[DEBUG] Validating auth token for API call")
            token_data = json.loads(base64.b64decode(auth_token).decode('utf-8'))

            # 檢查過期時間
            if token_data.get('exp', 0) < time.time():
                print(f"[DEBUG] API token expired")
                return jsonify({"error": "Authentication token expired"}), 403

            # 檢查角色
            if token_data.get('role') != 'Administrator':
                print(f"[DEBUG] API token invalid role: {token_data.get('role')}")
                return jsonify({"error": "Forbidden"}), 403

            # Token 有效，允許訪問
            print("[DEBUG] API token valid")
            return

        except Exception as e:
            print(f"[DEBUG] API token validation error: {str(e)}")
            return jsonify({"error": "Invalid authentication token", "details": str(e)}), 403

    # 如果沒有 token，檢查 Cookie
    cookie = request.headers.get('Cookie')
    if not cookie:
        print("[DEBUG] No auth token or cookie found")
        return jsonify({"error": "Unauthorized"}), 403

    try:
        # 驗證 Cookie
        print(f"[DEBUG] Validating cookie: {cookie[:100]}...")
        response = requests.get(ADMIN_API_VALIDATE, headers={"Cookie": cookie})
        if response.status_code != 200 or response.json().get("role") != "Administrator":
            print(f"[DEBUG] Cookie validation failed: status={response.status_code}, role={response.json().get('role')}")
            return jsonify({"error": "Forbidden"}), 403
        print("[DEBUG] Cookie validation successful")
    except Exception as e:
        print(f"[DEBUG] Cookie validation error: {str(e)}")
        return jsonify({"error": "Authentication failed", "details": str(e)}), 403

if __name__ == '__main__':
    port = int(os.environ.get('PORT', 5050))  # 改用 5050 端口
    print("🚀 Wiki.js 批量文檔導入服務啟動中...")
    print("📚 支援格式:", list(processor.supported_formats.keys()))
    print(f"🌐 服務地址: http://localhost:{port}")

    app.run(debug=True, host='0.0.0.0', port=port)
